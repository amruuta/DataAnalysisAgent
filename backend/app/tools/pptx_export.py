import json
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import pandas as pd
import structlog
from langchain_core.tools import tool

from app.config import settings

logger = structlog.get_logger(__name__)

_SUPPORTED_CHART_TYPES = {"bar", "line", "pie"}
_MAX_SLIDES = 12
_MAX_TABLE_ROWS = 12
_MAX_CHART_ROWS = 40


@dataclass(frozen=True)
class PresentationTheme:
    """Small visual system for generated decks."""

    navy: tuple[int, int, int] = (31, 41, 55)
    amber: tuple[int, int, int] = (245, 158, 11)
    slate: tuple[int, int, int] = (71, 85, 105)
    light: tuple[int, int, int] = (248, 250, 252)
    white: tuple[int, int, int] = (255, 255, 255)


def _strip_code_fence(value: str) -> str:
    text = value.strip()
    if not text.startswith("```"):
        return text

    lines = text.splitlines()
    if lines and lines[0].startswith("```"):
        lines = lines[1:]
    if lines and lines[-1].strip() == "```":
        lines = lines[:-1]
    text = "\n".join(lines)
    return text.strip()


def _safe_file_stem(value: str | None) -> str:
    stem = Path(value or "presentation").stem.strip().lower()
    cleaned = "".join(char if char.isalnum() or char in {"_", "-"} else "_" for char in stem)
    return cleaned.strip("_") or "presentation"


def _read_sql(sql_query: str, db_engine, max_rows: int) -> pd.DataFrame:
    """Read planned data with pandas and keep only the rows the slide needs."""
    df = pd.read_sql(sql_query, db_engine)
    return df.head(max(max_rows, 1))


def _parse_plan(plan_json: str) -> dict[str, Any]:
    try:
        plan = json.loads(_strip_code_fence(plan_json))
    except json.JSONDecodeError as exc:
        raise ValueError(f"Invalid JSON plan: {exc}") from exc

    if not isinstance(plan, dict):
        raise ValueError("The presentation plan must be a JSON object.")

    slides = plan.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("The presentation plan must include a non-empty slides list.")

    plan["slides"] = slides[:_MAX_SLIDES]
    return plan


class PptxDeckBuilder:
    """Builds a PowerPoint deck from a bounded JSON plan."""

    def __init__(self, db_engine, export_dir: str | Path):
        self.db_engine = db_engine
        self.export_dir = Path(export_dir)
        self.theme = PresentationTheme()

    def build(self, plan: dict[str, Any]) -> Path:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is not installed. Install backend requirements first."
            ) from exc

        self.export_dir.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, plan)
        for slide_plan in plan["slides"]:
            self._add_content_slide(prs, slide_plan)

        out_name = f"{uuid.uuid4().hex}_{_safe_file_stem(plan.get('file_name'))}.pptx"
        out_path = self.export_dir / out_name
        prs.save(out_path)
        return out_path

    def _add_title_slide(self, prs, plan: dict[str, Any]) -> None:
        from pptx.util import Inches, Pt

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._paint_background(slide, prs)
        self._add_accent_bar(slide, prs)

        title = str(plan.get("title") or "Data Presentation")
        subtitle = str(plan.get("subtitle") or "Generated from your data")

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.5), Inches(1.3))
        title_frame = title_box.text_frame
        title_frame.clear()
        paragraph = title_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(38)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self._rgb(self.theme.white)

        subtitle_box = slide.shapes.add_textbox(Inches(0.85), Inches(3.0), Inches(7.5), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        paragraph = subtitle_frame.paragraphs[0]
        paragraph.text = subtitle
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = self._rgb((226, 232, 240))

    def _add_content_slide(self, prs, slide_plan: dict[str, Any]) -> None:
        from pptx.util import Inches

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._paint_background(slide, prs, light=True)
        self._add_slide_title(slide, str(slide_plan.get("title") or "Insight"))

        content_top = Inches(1.28)
        left = Inches(0.72)
        width = Inches(5.25)

        self._add_bullets(slide, slide_plan.get("bullets", []), left, content_top, width)

        if isinstance(slide_plan.get("chart"), dict):
            self._add_chart(slide, slide_plan["chart"])
        elif isinstance(slide_plan.get("table"), dict):
            self._add_table(slide, slide_plan["table"])
        else:
            self._add_takeaway(slide, str(slide_plan.get("takeaway") or ""))

    def _add_chart(self, slide, chart_plan: dict[str, Any]) -> None:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.util import Inches

        chart_type = str(chart_plan.get("chart_type") or "bar").lower()
        if chart_type not in _SUPPORTED_CHART_TYPES:
            chart_type = "bar"

        max_rows = int(chart_plan.get("max_rows") or _MAX_CHART_ROWS)
        df = _read_sql(
            str(chart_plan["sql_query"]),
            self.db_engine,
            min(max_rows, _MAX_CHART_ROWS),
        )
        if df.empty:
            self._add_takeaway(slide, "No rows were returned for the planned chart.")
            return

        x_column = str(chart_plan.get("x_column") or df.columns[0])
        y_column = str(chart_plan.get("y_column") or df.columns[-1])
        if x_column not in df.columns or y_column not in df.columns:
            self._add_takeaway(slide, "The planned chart columns were not present in the query result.")
            return

        values = pd.to_numeric(df[y_column], errors="coerce").fillna(0).tolist()
        labels = [str(value) for value in df[x_column].tolist()]

        data = CategoryChartData()
        data.categories = labels
        data.add_series(str(y_column), values)

        chart_kind = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }[chart_type]

        chart = slide.shapes.add_chart(
            chart_kind,
            Inches(6.3),
            Inches(1.45),
            Inches(6.25),
            Inches(4.95),
            data,
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = str(chart_plan.get("title") or "Chart")
        chart.has_legend = chart_type == "pie"
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False

    def _add_table(self, slide, table_plan: dict[str, Any]) -> None:
        from pptx.util import Inches, Pt

        max_rows = int(table_plan.get("max_rows") or _MAX_TABLE_ROWS)
        df = _read_sql(
            str(table_plan["sql_query"]),
            self.db_engine,
            min(max_rows, _MAX_TABLE_ROWS),
        )
        if df.empty:
            self._add_takeaway(slide, "No rows were returned for the planned table.")
            return

        df = df.iloc[:, :5]
        rows = len(df) + 1
        cols = len(df.columns)
        table_shape = slide.shapes.add_table(
            rows,
            cols,
            Inches(6.25),
            Inches(1.45),
            Inches(6.35),
            Inches(4.8),
        )
        table = table_shape.table

        for column_index, column_name in enumerate(df.columns):
            cell = table.cell(0, column_index)
            cell.text = str(column_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._rgb(self.theme.navy)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = self._rgb(self.theme.white)

        for row_index, (_, row) in enumerate(df.iterrows(), start=1):
            for column_index, value in enumerate(row):
                cell = table.cell(row_index, column_index)
                cell.text = "" if pd.isna(value) else str(value)[:60]
                cell.text_frame.paragraphs[0].font.size = Pt(9)

    def _add_bullets(self, slide, bullets: Any, left, top, width) -> None:
        from pptx.util import Inches, Pt

        bullet_items = bullets if isinstance(bullets, list) else []
        bullet_items = [str(item) for item in bullet_items[:5] if str(item).strip()]
        if not bullet_items:
            bullet_items = ["Key context and findings from the selected data."]

        shape = slide.shapes.add_textbox(left, top, width, Inches(4.8))
        frame = shape.text_frame
        frame.clear()
        for index, item in enumerate(bullet_items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = self._rgb(self.theme.slate)

    def _add_takeaway(self, slide, text: str) -> None:
        from pptx.util import Inches, Pt

        message = text or "Use this slide for the main recommendation or next step."
        shape = slide.shapes.add_textbox(Inches(6.35), Inches(1.75), Inches(5.85), Inches(3.8))
        frame = shape.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = message
        paragraph.font.size = Pt(22)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self._rgb(self.theme.navy)

    def _add_slide_title(self, slide, title: str) -> None:
        from pptx.util import Inches, Pt

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.55))
        frame = title_box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(25)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self._rgb(self.theme.navy)

    def _paint_background(self, slide, prs, light: bool = False) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        fill_color = self.theme.light if light else self.theme.navy
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            prs.slide_height,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(fill_color)
        bg.line.fill.background()

        if light:
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(13.333),
                Inches(0.12),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = self._rgb(self.theme.amber)
            accent.line.fill.background()

    def _add_accent_bar(self, slide, prs) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(6.8),
            prs.slide_width,
            Inches(0.28),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(self.theme.amber)
        bar.line.fill.background()

    @staticmethod
    def _rgb(color: tuple[int, int, int]):
        from pptx.dml.color import RGBColor

        return RGBColor(*color)


def create_pptx_export_tool(db_engine):
    """Create a PPTX export tool that turns a JSON plan into a deck."""

    @tool
    def create_pptx_from_plan(presentation_plan_json: str) -> str:
        """Create a PowerPoint presentation from a JSON presentation plan.

        Use this after planning the deck content as JSON. The JSON must include:
        title, subtitle, file_name, and slides. Each slide can contain bullets and
        either a chart object or table object. Chart/table objects must include a
        SQL SELECT query and the columns needed for the visual.

        Returns:
            A markdown link to the generated PPTX file.
        """
        try:
            plan = _parse_plan(presentation_plan_json)
            out_path = PptxDeckBuilder(db_engine, settings.EXPORT_DIR).build(plan)
        except Exception as exc:
            logger.error("pptx_generation_failed", error=str(exc))
            return f"PPTX generation failed: {exc}"

        file_name = out_path.name
        export_link = f"{settings.PUBLIC_BASE_URL.rstrip('/')}/exports/{file_name}"
        logger.info("pptx_created", path=out_path, link=export_link)
        return f"PPTX file created: [{file_name}]({export_link})"

    return create_pptx_from_plan
